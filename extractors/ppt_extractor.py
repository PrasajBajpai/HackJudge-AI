from pptx import Presentation

def extract_text_from_ppt(ppt_path):

    presentation = Presentation(ppt_path)

    all_text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                all_text.append(shape.text)

    return "\n".join(all_text)